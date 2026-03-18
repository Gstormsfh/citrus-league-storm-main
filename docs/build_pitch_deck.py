#!/usr/bin/env python3
"""
Citrus Fantasy Sports — Investor Pitch Deck Generator
Matches the exact style of the existing deck (sage/dark green, orange accents, serif headings).
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ── Brand Colors ──
SAGE_BG = RGBColor(0xB8, 0xC3, 0xAD)
DARK_GREEN_BG = RGBColor(0x3A, 0x53, 0x35)
BRAND_GREEN = RGBColor(0x4A, 0x67, 0x41)
BRAND_ORANGE = RGBColor(0xD4, 0x70, 0x2A)
CREAM = RGBColor(0xF5, 0xF0, 0xE8)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x20, 0x20, 0x20)
LIGHT_ORANGE = RGBColor(0xFF, 0xE8, 0xCC)

# ── Fonts ──
HEADING_FONT = "Georgia"
BODY_FONT = "Calibri"


def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_textbox(slide, left, top, width, height, text, font_size=12,
                bold=False, italic=False, color=BLACK, font_name=BODY_FONT,
                alignment=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.italic = italic
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_multiline_textbox(slide, left, top, width, height, lines, default_size=11,
                          default_color=BLACK, font_name=BODY_FONT, alignment=PP_ALIGN.LEFT):
    """lines: list of (text, font_size, bold, italic, color) tuples"""
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    for i, line_data in enumerate(lines):
        if isinstance(line_data, str):
            text, size, bold, italic, clr = line_data, default_size, False, False, default_color
        else:
            text = line_data[0]
            size = line_data[1] if len(line_data) > 1 else default_size
            bold = line_data[2] if len(line_data) > 2 else False
            italic = line_data[3] if len(line_data) > 3 else False
            clr = line_data[4] if len(line_data) > 4 else default_color
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = text
        p.font.size = Pt(size)
        p.font.bold = bold
        p.font.italic = italic
        p.font.color.rgb = clr
        p.font.name = font_name
        p.alignment = alignment
        p.space_after = Pt(4)
    return txBox


def add_orange_circle(slide, left, top, size, text):
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(left), Inches(top), Inches(size), Inches(size))
    shape.fill.solid()
    shape.fill.fore_color.rgb = BRAND_ORANGE
    shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.font.name = HEADING_FONT
    p.alignment = PP_ALIGN.CENTER
    tf.paragraphs[0].space_before = Pt(0)
    return shape


def add_accent_bar(slide, left, top, width, height=0.06, color=BRAND_ORANGE):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_green_circle(slide, left, top, size=0.35):
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(left), Inches(top), Inches(size), Inches(size))
    shape.fill.solid()
    shape.fill.fore_color.rgb = BRAND_GREEN
    shape.line.fill.background()
    return shape


def build_deck():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]  # blank

    # ══════════════════════════════════════════════════════════════════
    # SLIDE 1: TITLE
    # ══════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, SAGE_BG)

    add_textbox(slide, 0.4, 0.3, 4, 0.4, "PRE-SEED PITCH  ·  ALPHA APPLICATION",
                font_size=9, color=BRAND_GREEN, font_name=BODY_FONT)

    add_textbox(slide, 0.4, 1.2, 5, 1.8, "CITRUS FANTASY\nSPORTS",
                font_size=48, bold=True, color=BLACK, font_name=HEADING_FONT)

    add_accent_bar(slide, 0.4, 3.3, 4.5)

    add_textbox(slide, 0.4, 3.6, 6, 0.6,
                "Building the world's most accurate Fantasy Hockey Platform",
                font_size=16, italic=False, color=BLACK, font_name=HEADING_FONT)

    add_multiline_textbox(slide, 5.5, 4.8, 4.2, 1.0, [
        ("ML-powered fantasy hockey built from 863K real NHL shots across 7 seasons.", 10, False, False, BLACK),
        ("31 engineered features · XGBoost v3 + Bayesian Regression · Claude AI (Stormy)", 10, False, False, BLACK),
    ])

    add_textbox(slide, 7.5, 6.8, 2.3, 0.4, "March 2026  |  Confidential",
                font_size=9, color=BRAND_GREEN)

    # ══════════════════════════════════════════════════════════════════
    # SLIDE 2: THE FOUNDER
    # ══════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, SAGE_BG)

    add_textbox(slide, 0.5, 0.4, 6, 1.2, "THE FOUNDER",
                font_size=44, bold=True, color=BLACK, font_name=HEADING_FONT)

    # Three cards
    cards = [
        ("01", "Professional Hockey\nPlayer",
         "Played at the professional level and coached at the elite level. Deep, lived understanding of hockey analytics \u2014 what matters on the ice."),
        ("02", "Chartered\nProfessional\nAccountant",
         "2.5 years of financial rigour. Understands unit economics, CAC, LTV, and how to build a sustainable business \u2014 not just a cool product."),
        ("03", "Self-Taught Data\nScientist",
         "Built the 70.9% R\u00b2 ML engine from scratch. XGBoost v3, Bayesian regression, GAR, GSAx. 95K+ lines of code. No co-founder needed."),
    ]

    for i, (num, title, desc) in enumerate(cards):
        x = 0.5 + i * 3.1
        # Card background
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(2.0), Inches(2.8), Inches(2.8))
        shape.fill.solid()
        shape.fill.fore_color.rgb = WHITE
        shape.line.color.rgb = BRAND_GREEN
        shape.line.width = Pt(1)

        add_orange_circle(slide, x + 0.1, 2.15, 0.45, num)
        add_textbox(slide, x + 0.7, 2.15, 2.0, 0.8, title,
                    font_size=13, bold=True, color=BLACK, font_name=HEADING_FONT)
        add_textbox(slide, x + 0.15, 3.1, 2.5, 1.5, desc,
                    font_size=9, color=BLACK)

    # The Journey
    add_textbox(slide, 0.5, 5.2, 5, 0.5, "THE JOURNEY",
                font_size=16, bold=True, color=BLACK, font_name=HEADING_FONT)

    add_textbox(slide, 0.5, 5.7, 9, 0.4,
                "2 years building nights and weekends while working full-time as a CPA.",
                font_size=11, color=BLACK)

    # Boxed italic text
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(6.2), Inches(9), Inches(0.8))
    shape.fill.solid()
    shape.fill.fore_color.rgb = WHITE
    shape.line.color.rgb = BRAND_ORANGE
    shape.line.width = Pt(2)
    add_textbox(slide, 0.7, 6.3, 8.6, 0.6,
                "Built the full-stack infrastructure, projection models, and product solo. 95K+ lines of code, 7 ML models, 60+ database tables \u2014 every line written by one person with a laptop and a vision.",
                font_size=10, italic=True, color=BLACK)

    # ══════════════════════════════════════════════════════════════════
    # SLIDE 3: WHAT IS FANTASY SPORTS?
    # ══════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, DARK_GREEN_BG)

    add_textbox(slide, 0.5, 0.3, 8, 1.2, "WHAT IS FANTASY SPORTS?",
                font_size=40, bold=True, color=WHITE, font_name=HEADING_FONT)
    add_textbox(slide, 0.5, 1.3, 6, 0.4, "Season Long Fantasy Leagues vs. Daily Fantasy Leagues",
                font_size=12, italic=True, color=CREAM)

    # Three phase cards
    phases = [
        ("DRAFT", [
            "Create a league with friends or join one publicly",
            "Draft real NHL players onto your own team",
            "Pick up, drop, and trade players with your league throughout the season",
        ]),
        ("MATCHUPS", [
            "Compete head-to-head against league mates each week",
            "Your players earn points based on real NHL performance",
            "Manage your lineup, make trades, and work the waiver wire",
        ]),
        ("CHAMPIONSHIP", [
            "After the regular season, top teams enter the playoffs",
            "Custom playoff formats crown a league champion",
            "Bragging rights, prizes, and a reason to watch every game",
        ]),
    ]

    for i, (title, bullets) in enumerate(phases):
        x = 0.5 + i * 3.1
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(2.0), Inches(2.8), Inches(3.0))
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(0x4A, 0x67, 0x41)
        shape.line.fill.background()

        add_textbox(slide, x + 0.15, 2.1, 2.5, 0.4, title,
                    font_size=16, bold=True, color=WHITE, font_name=HEADING_FONT)
        add_accent_bar(slide, x + 0.15, 2.55, 1.0, 0.04)

        bullet_text = "\n".join([f"\u2022  {b}" for b in bullets])
        add_textbox(slide, x + 0.15, 2.7, 2.5, 2.2, bullet_text,
                    font_size=9, color=CREAM)

    add_multiline_textbox(slide, 0.5, 5.5, 9, 1.5, [
        ("Citrus offers completely custom scoring and league styles", 14, True, False, WHITE),
        ("Opens the door for ALL users and ALL league types \u2014 from casual to hardcore", 11, False, False, CREAM),
    ], font_name=HEADING_FONT)

    # ══════════════════════════════════════════════════════════════════
    # SLIDE 4: THE PROBLEM
    # ══════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, SAGE_BG)

    add_textbox(slide, 0.5, 0.3, 6, 1.0, "THE PROBLEM",
                font_size=44, bold=True, color=BLACK, font_name=HEADING_FONT)
    add_textbox(slide, 0.5, 1.2, 9, 0.5,
                "NHL is the fastest-growing Big 4 sport \u2014 and the only one abandoned by major fantasy platforms.",
                font_size=12, italic=True, color=BLACK)

    # Big stat
    add_textbox(slide, 0.5, 2.2, 2.5, 0.8, "5.8M",
                font_size=48, bold=True, color=BRAND_ORANGE, font_name=HEADING_FONT)
    add_textbox(slide, 0.5, 3.0, 3, 0.6,
                "Adults playing fantasy hockey at the end of 2025 \u2014 Mordor Research",
                font_size=10, color=BLACK)

    # Problem cards
    problems = [
        ("Stale, Inaccurate Projections",
         "Yahoo and ESPN use basic averages that haven't evolved in years. No ML, no play coordinate data, no xG. Projections regularly miss by 30-50%."),
        ("No Analytics Depth",
         "Advanced metrics like xG, xA, GAR, and GSAx don't exist on any consumer fantasy platform. Serious players have nowhere to go."),
        ("High Barrier to Entry",
         "New players feel overwhelmed. Experienced players feel underserved. No platform bridges both worlds."),
    ]

    for i, (title, desc) in enumerate(problems):
        y = 2.2 + i * 1.3
        add_accent_bar(slide, 4.0, y + 0.15, 0.08, 0.7, BRAND_ORANGE)
        add_textbox(slide, 4.3, y, 5.2, 0.4, title,
                    font_size=14, bold=True, color=BLACK, font_name=HEADING_FONT)
        add_textbox(slide, 4.3, y + 0.4, 5.2, 0.7, desc,
                    font_size=10, color=BLACK)

    add_textbox(slide, 0.5, 6.5, 9, 0.5,
                "NHL viewership up 35% since 2020  ·  14.6% CAGR \u2014 fastest-growing Big 4  ·  2026 Winter Olympics (Milan) featuring hockey",
                font_size=9, italic=True, color=BRAND_GREEN)

    # ══════════════════════════════════════════════════════════════════
    # SLIDE 5: THE SOLUTION
    # ══════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, DARK_GREEN_BG)

    add_textbox(slide, 0.5, 0.3, 6, 1.0, "THE SOLUTION",
                font_size=44, bold=True, color=WHITE, font_name=HEADING_FONT)
    add_textbox(slide, 0.5, 1.2, 6, 0.4, "One platform. Four pillars. Every advantage.",
                font_size=13, italic=True, color=CREAM)

    pillars = [
        ("CORE ENGINE", "ML Projections", "0.817 AUC-ROC",
         "If you show our model a random goal and a random save, it correctly identifies which is which 81.7% of the time. Born from 8 NHL seasons of data, and 870,000 shots analyzed"),
        ("FAN ENGAGEMENT", "Stormy AI", "Claude-Powered",
         "Your personal AI General Manager. Ask lineup questions, evaluate trades, get waiver picks \u2014 in plain English. Real-time, roster-aware advice."),
        ("LIVE EXPERIENCE", "Real-Time Scoring", "",
         "Live WebSocket scoring that updates every 30 seconds during games. Sub-second response. No refreshing, no waiting."),
        ("ANALYTICS LAYER", "Advanced Analytics", "xG \xb7 xA \xb7 GAR \xb7 GSAx",
         "Professional-grade metrics on a consumer platform for the first time. 7 ML models powering the same data NHL front offices use."),
    ]

    for i, (label, title, badge, desc) in enumerate(pillars):
        col = i % 2
        row = i // 2
        x = 0.5 + col * 4.7
        y = 2.0 + row * 2.6

        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(4.3), Inches(2.3))
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(0x4A, 0x67, 0x41)
        shape.line.fill.background()

        add_textbox(slide, x + 0.15, y + 0.1, 1.5, 0.3, label,
                    font_size=8, color=CREAM)
        add_textbox(slide, x + 0.15, y + 0.35, 2.5, 0.4, title,
                    font_size=18, bold=True, color=WHITE, font_name=HEADING_FONT)
        if badge:
            add_textbox(slide, x + 2.5, y + 0.1, 1.6, 0.4, badge,
                        font_size=10, bold=True, color=BRAND_ORANGE, font_name=HEADING_FONT,
                        alignment=PP_ALIGN.RIGHT)
        add_textbox(slide, x + 0.15, y + 0.85, 4.0, 1.3, desc,
                    font_size=9, color=CREAM)

    # ══════════════════════════════════════════════════════════════════
    # SLIDE 6: BEST-IN-CLASS PROJECTIONS
    # ══════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, SAGE_BG)

    add_textbox(slide, 0.5, 0.3, 9, 1.0, "BEST-IN-CLASS PROJECTIONS",
                font_size=40, bold=True, color=BLACK, font_name=HEADING_FONT)
    add_textbox(slide, 0.5, 1.2, 6, 0.4,
                "Real season data  \xb7  Citrus vs Yahoo vs ESPN vs Actual",
                font_size=12, italic=True, color=BLACK)

    # Comparison table
    table_data = [
        ("Player", "Citrus", "Yahoo", "ESPN", "Actual"),
        ("Dylan Larkin", "26.0", "14.7", "16.0", "24"),
        ("Emil Heineman", "15.8", "12.8", "11.3", "15"),
        ("Vince Dunn", "6.2", "5.2", "4.6", "7"),
        ("Josh Doan", "19.2", "13.2", "9.6", "18"),
    ]

    from pptx.util import Inches as In
    table = slide.shapes.add_table(5, 5, In(0.5), In(2.0), In(5.5), In(2.5)).table
    table.columns[0].width = In(1.5)
    for i in range(1, 5):
        table.columns[i].width = In(1.0)

    for r, row_data in enumerate(table_data):
        for ci, val in enumerate(row_data):
            cell = table.cell(r, ci)
            cell.text = val
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(11)
            p.font.name = HEADING_FONT
            if r == 0:
                p.font.bold = True
                p.font.color.rgb = WHITE
                cell.fill.solid()
                cell.fill.fore_color.rgb = BRAND_GREEN
            else:
                p.font.color.rgb = BLACK
                if ci == 1:  # Citrus column
                    p.font.bold = True
                    p.font.color.rgb = BRAND_GREEN
                elif ci == 4:  # Actual column
                    p.font.bold = True
                    p.font.color.rgb = BRAND_ORANGE

    # 70.9% callout
    add_textbox(slide, 6.5, 2.0, 3, 0.8, "70.9%",
                font_size=48, bold=True, color=BRAND_ORANGE, font_name=HEADING_FONT)
    add_multiline_textbox(slide, 6.5, 2.8, 3, 0.8, [
        ("Player-Season R\u00b2", 14, True, False, BLACK),
        ("0.817 AUC-ROC  \xb7  0.842 Correlation", 9, False, False, BLACK),
        ("591 players validated  \xb7  863K shots", 9, False, False, BLACK),
    ], font_name=HEADING_FONT)

    # MOAT section
    add_textbox(slide, 0.5, 4.8, 5, 0.5, "THE MOAT",
                font_size=18, bold=True, color=BLACK, font_name=HEADING_FONT)
    add_textbox(slide, 0.5, 5.25, 5, 0.4,
                "Can't Yahoo or ESPN just copy this?  No.",
                font_size=13, bold=True, color=BRAND_ORANGE, font_name=HEADING_FONT)
    add_textbox(slide, 0.5, 5.7, 9, 0.8,
                "2+ years of domain-specific engineering across 7 NHL seasons and 863K shots. Our v3 model now exceeds MoneyPuck's published AUC-ROC \u2014 with a proprietary calibration pipeline and pass context system (AUC 0.961) that no competitor replicates. This isn't a weekend project. It's a moat.",
                font_size=10, color=BLACK)

    add_textbox(slide, 0.5, 6.6, 9, 0.5,
                "XGBoost v3 on exact play coordinates \xb7 Proprietary multi-layer calibration \xb7 Bayesian per-game projections \xb7 Opponent-aware difficulty rating",
                font_size=9, italic=True, color=BRAND_GREEN)

    # ══════════════════════════════════════════════════════════════════
    # SLIDE 7: MARKET OPPORTUNITY
    # ══════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, SAGE_BG)

    add_textbox(slide, 0.5, 0.3, 8, 1.0, "MARKET OPPORTUNITY",
                font_size=40, bold=True, color=BLACK, font_name=HEADING_FONT)
    add_textbox(slide, 0.5, 1.2, 6, 0.4, "Global audience  \xb7  Growing sector  \xb7  Clear path to scale",
                font_size=12, italic=True, color=BLACK)

    # Big stats
    add_textbox(slide, 0.5, 2.0, 2, 0.8, "$32B+",
                font_size=44, bold=True, color=BRAND_ORANGE, font_name=HEADING_FONT)
    add_textbox(slide, 0.5, 2.7, 2, 0.5, "Fantasy Sports TAM\n2025 global est.",
                font_size=10, color=BLACK)

    add_textbox(slide, 2.8, 2.0, 2, 0.8, "1-2B",
                font_size=36, bold=True, color=BRAND_GREEN, font_name=HEADING_FONT)
    add_textbox(slide, 2.8, 2.7, 2, 0.5, "Fantasy Hockey SAM\n2025 global est.",
                font_size=10, color=BLACK)

    add_textbox(slide, 5.0, 2.0, 2, 0.8, "14.6%",
                font_size=44, bold=True, color=BRAND_ORANGE, font_name=HEADING_FONT)
    add_textbox(slide, 5.0, 2.7, 2, 0.5, "CAGR \u2014 Fastest\nGrowing Big 4",
                font_size=10, color=BLACK)

    add_textbox(slide, 7.5, 2.0, 2, 0.8, "5.8M",
                font_size=36, bold=True, color=BRAND_GREEN, font_name=HEADING_FONT)
    add_textbox(slide, 7.5, 2.7, 2, 0.5, "Active Fantasy\nHockey Adults (NA)",
                font_size=10, color=BLACK)

    # Global audience
    add_textbox(slide, 0.5, 3.5, 5, 0.4, "A GLOBAL AUDIENCE FOR THE NHL",
                font_size=13, bold=True, color=BLACK, font_name=HEADING_FONT)
    add_textbox(slide, 0.5, 3.9, 5, 0.4,
                "Hockey fans worldwide follow the NHL \u2014 but have zero dedicated fantasy tools.",
                font_size=10, italic=True, color=BLACK)
    global_points = [
        "\u25b8 32% of NHL players born outside North America",
        "\u25b8 20 birth countries on 2025-26 NHL rosters",
        "\u25b8 IIHF Worlds: broadcast in 100+ territories",
        "\u25b8 50M+ hockey fans across NA & Europe",
        "\u25b8 84 IIHF member nations \u2014 untapped audience",
        "\u25b8 Europe fantasy market: $6.9B, 12.8% CAGR",
    ]
    add_textbox(slide, 0.5, 4.4, 5, 2.0, "\n".join(global_points),
                font_size=9, color=BLACK)

    # Competitive landscape
    add_textbox(slide, 5.5, 3.5, 4, 0.4, "COMPETITIVE LANDSCAPE",
                font_size=13, bold=True, color=BLACK, font_name=HEADING_FONT)

    competitors = [
        ("Yahoo Fantasy", "No NHL \xb7 Stale projections \xb7 Declining", "~60%", BRAND_ORANGE),
        ("ESPN Fantasy", "No NHL \xb7 Same bad models \xb7 Losing users", "~20%", BRAND_ORANGE),
        ("FanTrax / Others", "Niche \xb7 No scale \xb7 No ML tools", "~20%", BRAND_ORANGE),
        ("Citrus Fantasy", "Best-in-class ML \xb7 AI-powered \xb7 Mar 2026 launch", "0% \u2192 \u221e", BRAND_GREEN),
    ]

    for i, (name, desc, share, clr) in enumerate(competitors):
        y = 4.0 + i * 0.75
        add_accent_bar(slide, 5.5, y + 0.35, 3.5, 0.04, clr)
        add_textbox(slide, 5.5, y, 3, 0.3, name,
                    font_size=12, bold=True, color=BLACK, font_name=HEADING_FONT)
        add_textbox(slide, 5.5, y + 0.25, 3, 0.3, desc,
                    font_size=8, color=BLACK)
        add_textbox(slide, 8.2, y, 1.2, 0.3, share,
                    font_size=16, bold=True, color=clr, font_name=HEADING_FONT,
                    alignment=PP_ALIGN.RIGHT)

    # ══════════════════════════════════════════════════════════════════
    # SLIDE 8: EARLY TRACTION (NEW)
    # ══════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, DARK_GREEN_BG)

    add_textbox(slide, 0.5, 0.3, 8, 1.0, "EARLY TRACTION",
                font_size=44, bold=True, color=WHITE, font_name=HEADING_FONT)
    add_textbox(slide, 0.5, 1.2, 7, 0.4,
                "Live product  \xb7  Real users  \xb7  Growing engagement  \xb7  Zero marketing spend",
                font_size=13, italic=True, color=CREAM)

    # Four big metric cards
    metrics = [
        ("4", "Active Leagues", "Multiple league formats running\nconcurrently in-season"),
        ("50+", "Active Users", "Organic growth through word-of-mouth\nand hockey communities"),
        ("100+", "Daily Visits", "Consistent daily engagement\nacross the user base"),
        ("2+ min", "Avg Engagement", "Deep sessions \u2014 users aren't\nbouncing, they're managing"),
    ]

    for i, (big_num, label, desc) in enumerate(metrics):
        col = i % 2
        row = i // 2
        x = 0.5 + col * 4.7
        y = 2.2 + row * 2.4

        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(4.3), Inches(2.1))
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(0x4A, 0x67, 0x41)
        shape.line.fill.background()

        add_textbox(slide, x + 0.3, y + 0.2, 2.5, 0.8, big_num,
                    font_size=48, bold=True, color=BRAND_ORANGE, font_name=HEADING_FONT)
        add_textbox(slide, x + 0.3, y + 0.95, 3.5, 0.4, label,
                    font_size=16, bold=True, color=WHITE, font_name=HEADING_FONT)
        add_accent_bar(slide, x + 0.3, y + 1.35, 2.0, 0.04)
        add_textbox(slide, x + 0.3, y + 1.45, 3.5, 0.6, desc,
                    font_size=9, color=CREAM)

    add_textbox(slide, 0.5, 7.0, 9, 0.3,
                "Alpha launched during 2025-26 NHL season  \xb7  All growth organic  \xb7  Product-market fit signal with zero paid acquisition",
                font_size=9, italic=True, color=CREAM)

    # ══════════════════════════════════════════════════════════════════
    # SLIDE 9: BUSINESS MODEL + GTM
    # ══════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, SAGE_BG)

    add_textbox(slide, 0.5, 0.3, 8, 1.0, "BUSINESS MODEL + GTM",
                font_size=40, bold=True, color=BLACK, font_name=HEADING_FONT)
    add_textbox(slide, 0.5, 1.15, 6, 0.4,
                "Freemium launch  \xb7  Path to 100K users  \xb7  Global audience expansion",
                font_size=12, italic=True, color=BLACK)

    # Revenue model
    add_textbox(slide, 0.5, 1.8, 3, 0.4, "REVENUE MODEL",
                font_size=14, bold=True, color=BLACK, font_name=HEADING_FONT)

    tiers = [
        ("FREE TIER", "$0", "Basic projections \xb7 Standard scoring"),
        ("PREMIUM", "$4.99/mo", "Advanced ML (70.9% R\u00b2) \xb7 Full Stormy AI \xb7 Schedule optimizer"),
        ("DFS EXPANSION", "12% rake", "Daily Fantasy Sports at 100K users \xb7 Season-long + weekly"),
    ]

    for i, (name, price, desc) in enumerate(tiers):
        y = 2.3 + i * 0.9
        add_textbox(slide, 0.5, y, 2, 0.3, name,
                    font_size=13, bold=True, color=BLACK, font_name=HEADING_FONT)
        add_textbox(slide, 0.5, y + 0.3, 3.5, 0.3, desc, font_size=9, color=BLACK)
        add_textbox(slide, 3.2, y, 1.2, 0.3, price,
                    font_size=13, bold=True, color=BRAND_ORANGE, font_name=HEADING_FONT)

    # Path to scale
    add_textbox(slide, 5.0, 1.8, 4.5, 0.4, "PATH TO SCALE",
                font_size=14, bold=True, color=BLACK, font_name=HEADING_FONT)

    milestones = [
        ("Q1\u2013Q2 2026", "Product-market fit \xb7 Beta expansion \xb7 App Store launch",
         "1,000 users", "$5K"),
        ("Q3\u2013Q4 2026", "Premium launch \xb7 Projection accuracy marketing \xb7 Reddit/X",
         "10,000 users", "$72K"),
        ("Q1\u2013Q2 2027", "European NHL fan expansion \xb7 Influencer partnerships",
         "25,000 users", "$174K"),
        ("Q3\u2013Q4 2027", "DFS launch \xb7 Full vision \xb7 Revenue optimization",
         "100,000 users", "$1.3M"),
    ]

    for i, (period, desc, users, rev) in enumerate(milestones):
        y = 2.3 + i * 1.1
        add_green_circle(slide, 5.0, y + 0.05)
        add_textbox(slide, 5.5, y, 2.5, 0.25, period,
                    font_size=11, bold=True, color=BLACK, font_name=HEADING_FONT)
        add_textbox(slide, 5.5, y + 0.25, 3.0, 0.3, desc,
                    font_size=8, color=BLACK)
        add_textbox(slide, 8.0, y, 1.5, 0.25, users,
                    font_size=11, bold=True, color=BRAND_GREEN, font_name=HEADING_FONT,
                    alignment=PP_ALIGN.RIGHT)
        add_textbox(slide, 8.5, y + 0.25, 1.0, 0.25, rev,
                    font_size=11, bold=True, color=BRAND_ORANGE, font_name=HEADING_FONT,
                    alignment=PP_ALIGN.RIGHT)

    # Bottom comparison
    add_textbox(slide, 0.5, 5.6, 9, 0.4,
                "Sleeper (NFL Specific Comp) Grew to 1,000,000 users in their first year of operation in 2019",
                font_size=9, italic=True, color=BRAND_GREEN)
    add_textbox(slide, 0.5, 6.0, 9, 0.4,
                "Global expansion: European NHL fans (Sweden, Finland, Czech Republic, Germany) \u2192 Future: NBA / NFL cross-sport",
                font_size=9, bold=True, italic=False, color=BRAND_GREEN)

    # ══════════════════════════════════════════════════════════════════
    # SLIDE 10: THE ASK (UPDATED — Funding breakdown)
    # ══════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, DARK_GREEN_BG)

    add_textbox(slide, 0.5, 0.3, 6, 1.0, "THE ASK",
                font_size=44, bold=True, color=WHITE, font_name=HEADING_FONT)

    # Big number
    add_textbox(slide, 0.5, 1.3, 5, 0.8, "$250K PRE-SEED",
                font_size=36, bold=True, color=BRAND_ORANGE, font_name=HEADING_FONT)
    add_textbox(slide, 0.5, 2.0, 6, 0.4,
                "24-month runway to $1M+ gross revenue  \xb7  Cash-flow positive by Month 9",
                font_size=11, italic=True, color=CREAM)

    # Funding breakdown
    add_textbox(slide, 0.5, 2.7, 4, 0.4, "WHERE IT GOES",
                font_size=14, bold=True, color=WHITE, font_name=HEADING_FONT)
    add_accent_bar(slide, 0.5, 3.05, 4.0, 0.04)

    allocations = [
        ("Licensing, Legal & Compliance", "$80K", "32%",
         "Criminal Code opinion, gaming counsel retainer, contest rules, DFS regulatory filings"),
        ("Marketing & User Acquisition", "$65K", "26%",
         "Paid channels, influencer/content partnerships, Oct 2026 season push"),
        ("Team / Engineering Contractors", "$50K", "20%",
         "Pull hiring forward to Month 4. Accelerate product velocity."),
        ("Infrastructure & Scaling", "$25K", "10%",
         "Supabase Team tier, Claude API at scale, proxy rotation"),
        ("Operating Reserve", "$30K", "12%",
         "4+ months buffer at peak burn. Never let cash hit zero."),
    ]

    for i, (name, amount, pct, desc) in enumerate(allocations):
        y = 3.2 + i * 0.75
        add_orange_circle(slide, 0.5, y, 0.35, str(i + 1))
        add_textbox(slide, 1.0, y - 0.02, 3.0, 0.3, name,
                    font_size=11, bold=True, color=WHITE)
        add_textbox(slide, 4.0, y - 0.02, 0.8, 0.3, amount,
                    font_size=11, bold=True, color=BRAND_ORANGE,
                    alignment=PP_ALIGN.RIGHT)
        add_textbox(slide, 1.0, y + 0.28, 3.8, 0.3, desc,
                    font_size=8, color=CREAM)

    # Right side: Why invest
    add_textbox(slide, 5.5, 2.7, 4, 0.4, "WHY THIS IS A GOOD BET",
                font_size=14, bold=True, color=WHITE, font_name=HEADING_FONT)
    add_accent_bar(slide, 5.5, 3.05, 4.0, 0.04)

    reasons = [
        ("Capital efficient", "$250K in \u2192 $450K Y2 EBITDA = 1.8x return in 24 months"),
        ("Proven product", "Live alpha with 50+ users, 4 leagues, 2+ min engagement \u2014 zero spend"),
        ("Legally buttoned up", "$80K for gaming counsel, Criminal Code opinion, contest rules \u2014 no shortcuts"),
        ("Massive LTV:CAC", "$2.50 CAC vs DraftKings\u2019 $200+. LTV:CAC ratio hits 10x+ in-season"),
        ("97% gross margins", "Pure software. Near-zero marginal COGS at scale"),
        ("Solo founder = lean", "CPA + data scientist + hockey pro. No co-founder bloat"),
    ]

    for i, (title, desc) in enumerate(reasons):
        y = 3.2 + i * 0.62
        add_textbox(slide, 5.5, y, 4, 0.25, title,
                    font_size=11, bold=True, color=BRAND_ORANGE)
        add_textbox(slide, 5.5, y + 0.23, 4, 0.3, desc,
                    font_size=8, color=CREAM)

    add_textbox(slide, 0.5, 7.0, 9, 0.3,
                "citrusfantasy.com  \xb7  garrett@citrusfantasy.com  \xb7  Confidential",
                font_size=9, color=CREAM, alignment=PP_ALIGN.CENTER)

    # ══════════════════════════════════════════════════════════════════
    # SLIDE 11: APPENDIX DIVIDER
    # ══════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, WHITE)

    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(3.5), Inches(1.5), Inches(3), Inches(4.5))
    shape.fill.solid()
    shape.fill.fore_color.rgb = BRAND_GREEN
    shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.text = "APPENDIX"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.font.name = HEADING_FONT
    p.alignment = PP_ALIGN.CENTER

    # ══════════════════════════════════════════════════════════════════
    # SLIDE 12: TRACTION + THE TEAM (UPDATED metrics)
    # ══════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, DARK_GREEN_BG)

    add_textbox(slide, 0.5, 0.3, 8, 1.0, "TRACTION + THE TEAM",
                font_size=40, bold=True, color=WHITE, font_name=HEADING_FONT)
    add_textbox(slide, 0.5, 1.15, 6, 0.4,
                "Fully functional  \xb7  Live product  \xb7  Solo founder  \xb7  Undeniable background",
                font_size=12, italic=True, color=CREAM)

    # Traction metrics (UPDATED)
    traction = [
        ("2 Years", "Solo Build\nNights & Weekends"),
        ("50+", "Active Users\nIn-Season"),
        ("4", "Active Leagues\nRunning"),
        ("$1K", "Total Invested\nBootstrapped"),
        ("70.9%", "Player R\u00b2\nValidated"),
        ("Mar '26", "iOS App Store\nLaunch"),
    ]

    for i, (big, small) in enumerate(traction):
        col = i % 3
        row = i // 3
        x = 0.5 + col * 3.1
        y = 2.0 + row * 1.6

        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(2.8), Inches(1.3))
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(0x4A, 0x67, 0x41)
        shape.line.fill.background()

        clr = BRAND_ORANGE if i % 2 == 0 else WHITE
        add_textbox(slide, x + 0.15, y + 0.1, 2.5, 0.5, big,
                    font_size=28, bold=True, color=clr, font_name=HEADING_FONT)
        add_textbox(slide, x + 0.15, y + 0.65, 2.5, 0.5, small,
                    font_size=9, color=CREAM)

    # Founder section
    add_textbox(slide, 0.5, 5.5, 4, 0.4, "THE FOUNDER",
                font_size=14, bold=True, color=WHITE, font_name=HEADING_FONT)

    founder_cards = [
        ("01", "Professional Hockey Player",
         "Played at the professional level. Deep, lived understanding of what matters, what scouts track, what analytics the pros rely on."),
        ("02", "Chartered Professional Accountant",
         "2.5 years of financial rigour. Understands unit economics, CAC, LTV, and how to build a sustainable business."),
        ("03", "Self-Taught Data Scientist",
         "Built the 70.9% R\u00b2 ML engine from scratch. XGBoost, Bayesian regression, GAR, GSAx. 95K+ lines. No co-founder."),
    ]

    for i, (num, title, desc) in enumerate(founder_cards):
        x = 0.5 + i * 3.1
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(5.9), Inches(2.8), Inches(1.3))
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(0x4A, 0x67, 0x41)
        shape.line.fill.background()

        add_orange_circle(slide, x + 0.1, 6.0, 0.35, num)
        add_textbox(slide, x + 0.55, 6.0, 2.1, 0.35, title,
                    font_size=10, bold=True, color=WHITE, font_name=HEADING_FONT)
        add_textbox(slide, x + 0.15, 6.4, 2.5, 0.7, desc,
                    font_size=8, color=CREAM)

    # ══════════════════════════════════════════════════════════════════
    # SLIDE 13: HOW WE SEE THE ICE
    # ══════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, DARK_GREEN_BG)

    add_textbox(slide, 0.5, 0.3, 8, 1.0, "HOW WE SEE THE ICE",
                font_size=40, bold=True, color=WHITE, font_name=HEADING_FONT)
    add_textbox(slide, 0.5, 1.2, 6, 0.4,
                "Every shot  \xb7  31 engineered features  \xb7  863K shots trained",
                font_size=12, italic=True, color=CREAM)

    # xG zones visualization (simplified bar chart)
    zones = [
        ("Deep Zone\n0.05-0.1", 1.0, RGBColor(0x8B, 0x45, 0x13)),
        ("High Slot\n0.2-0.3", 2.5, RGBColor(0xA0, 0x8B, 0x40)),
        ("Slot\n0.3-0.4", 3.5, RGBColor(0x6B, 0x8E, 0x23)),
        ("Crease\n0.5-0.6", 4.5, RGBColor(0x8B, 0x00, 0x00)),
    ]

    for i, (label, bar_h, clr) in enumerate(zones):
        x = 1.0 + i * 1.8
        bar_top = 6.5 - bar_h
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                       Inches(x), Inches(bar_top), Inches(1.3), Inches(bar_h))
        shape.fill.solid()
        shape.fill.fore_color.rgb = clr
        shape.line.fill.background()
        add_textbox(slide, x, 6.6, 1.3, 0.6, label,
                    font_size=9, bold=True, color=CREAM, alignment=PP_ALIGN.CENTER)

    # Model stats
    add_textbox(slide, 5.5, 2.0, 4, 0.4, "xG v3 MODEL",
                font_size=16, bold=True, color=WHITE, font_name=HEADING_FONT)

    add_textbox(slide, 5.5, 2.5, 4, 1.5,
                "XGBClassifier trained on 863K real NHL shots across 7 seasons. Each shot gets 31 features extracted and a goal probability assigned before the puck crosses the line.",
                font_size=10, color=CREAM)

    model_stats = [
        ("AUC-ROC", "0.817", "Shot-level classification"),
        ("Brier Score", "0.054", "Lower = better (perfect: 0)"),
        ("Player R\u00b2", "70.9%", "Season-level variance"),
        ("Log Loss", "0.200", "Prediction uncertainty"),
    ]

    for i, (label, value, desc) in enumerate(model_stats):
        y = 4.0 + i * 0.65
        add_textbox(slide, 5.5, y, 1.5, 0.25, label,
                    font_size=9, color=CREAM)
        add_textbox(slide, 7.0, y, 1.2, 0.25, value,
                    font_size=18, bold=True, color=BRAND_ORANGE, font_name=HEADING_FONT)
        add_textbox(slide, 8.2, y, 1.5, 0.25, desc,
                    font_size=8, color=CREAM)

    add_textbox(slide, 5.5, 6.8, 4, 0.4,
                "Max xG capped at 0.60 per shot \xb7 Multi-layer calibration pipeline \xb7 7 seasons of NHL data",
                font_size=8, italic=True, color=CREAM)

    # ══════════════════════════════════════════════════════════════════
    # SLIDE 14: THE ENGINE
    # ══════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, SAGE_BG)

    add_textbox(slide, 0.5, 0.3, 6, 1.0, "THE ENGINE",
                font_size=44, bold=True, color=BLACK, font_name=HEADING_FONT)
    add_textbox(slide, 0.5, 1.2, 6, 0.4,
                "End-to-end proprietary stack  \xb7  Built solo  \xb7  2+ years",
                font_size=12, italic=True, color=BLACK)

    # Big stats
    engine_stats = [
        ("95K+", "Lines of Code"),
        ("7", "ML Models"),
        ("60+", "Database Tables"),
        ("30+", "xG Features"),
        ("30s", "Live Latency"),
    ]
    for i, (val, label) in enumerate(engine_stats):
        x = 0.5 + i * 1.8
        add_textbox(slide, x, 1.9, 1.5, 0.6, val,
                    font_size=32, bold=True, color=BRAND_ORANGE, font_name=HEADING_FONT)
        add_textbox(slide, x, 2.5, 1.5, 0.3, label,
                    font_size=10, color=BLACK)

    # Architecture pipeline
    pipeline = [
        ("1", "Data Ingestion", "4 NHL API endpoints\n100-IP proxy rotation\nCircuit breaker + backoff"),
        ("2", "PostgreSQL DB", "60+ tables \xb7 200+ migrations\n45+ RLS policies\nImmutable raw data"),
        ("3", "ML Engine", "", ""),
        ("4", "Projection Engine", "Bayesian shrinkage\nDDR + finishing talent", ""),
        ("5", "Live App + Stormy", "iOS launch Mar 2026", ""),
    ]

    for i, item in enumerate(pipeline):
        x = 0.5 + i * 1.9
        y = 3.5

        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(1.7), Inches(2.5))
        shape.fill.solid()
        shape.fill.fore_color.rgb = WHITE
        shape.line.color.rgb = BRAND_GREEN
        shape.line.width = Pt(1)

        add_orange_circle(slide, x + 0.6, y + 0.1, 0.45, item[0])
        add_textbox(slide, x + 0.1, y + 0.65, 1.5, 0.5, item[1],
                    font_size=11, bold=True, color=BLACK, font_name=HEADING_FONT,
                    alignment=PP_ALIGN.CENTER)
        if item[2]:
            add_textbox(slide, x + 0.1, y + 1.2, 1.5, 1.2, item[2],
                        font_size=8, color=BLACK, alignment=PP_ALIGN.CENTER)

    add_textbox(slide, 0.5, 6.5, 9, 0.5,
                "Enterprise-grade: 7-layer defence \xb7 45+ RLS policies \xb7 PKCE auth \xb7 Firebase Hosting \xb7 Supabase Realtime \xb7 100-IP proxy rotation",
                font_size=9, italic=True, color=BRAND_GREEN)

    # ── Save ──
    output_path = os.path.join(os.path.dirname(os.path.abspath(__file__)), "Citrus_Pitch_Deck.pptx")
    prs.save(output_path)
    return output_path


if __name__ == "__main__":
    path = build_deck()
    print(f"\nPitch deck saved to: {path}")
    print("14 slides including new Early Traction + updated The Ask")
